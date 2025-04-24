import requests
import json
import os
import uuid
from pptx import Presentation
from pptx.enum.text import PP_ALIGN


# Set up API credentials
subscription_key = ''  # Your Azure Translator subscription key
endpoint = 'https://api.cognitive.microsofttranslator.com/'  # The endpoint URL for the Translator API
location = ''  # Your resource location e.g eastus


# API setup
path = '/translate'
constructed_url = endpoint + path


params = {
   'api-version': '3.0',  # API version
   'from': 'en',  # Source language
   'to': ['ar']  # List of target languages (can add more if needed)
}


headers = {
   'Ocp-Apim-Subscription-Key': subscription_key,
   'Ocp-Apim-Subscription-Region': location,  # Required for regional resources
   'Content-type': 'application/json',
   'X-ClientTraceId': str(uuid.uuid4())  # Unique identifier for the request
}


def translate_text(text):
   if not text.strip():  # Skips any empty text on document
       return text


   body = [{'text': text}]
   response = requests.post(constructed_url, params=params, headers=headers, json=body)


   if response.status_code != 200:
       print(f"Error: Received status code {response.status_code}")
       print(f"Response: {response.text}")
       return text


   try:
       response_json = response.json()
       return response_json[0]['translations'][0]['text']
   except (KeyError, IndexError, json.JSONDecodeError) as e:
       print(f"Error parsing response: {e}")
       print(f"Response content: {response.content}")
       return text


def set_arabic_format(paragraph): # Sets the text to RTL writing direction
   paragraph.alignment = PP_ALIGN.RIGHT
   for run in paragraph.runs:
       run.font.name = 'Arial'
       run.font.rtl = True


def translate_presentation(input_file, output_file):
   prs = Presentation(input_file)
   for slide in prs.slides:  # Iterative scan through each slide in the presentation
       # Translate slide content
       for shape in slide.shapes:
           if shape.has_text_frame:
               for paragraph in shape.text_frame.paragraphs:
                   for run in paragraph.runs:
                       translated_text = translate_text(run.text)
                       run.text = translated_text
                   set_arabic_format(paragraph)
           elif shape.shape_type == 19:  # Check if the shape is a table
               table = shape.table
               for row in table.rows:
                   for cell in row.cells:
                       for paragraph in cell.text_frame.paragraphs:
                           for run in paragraph.runs:
                               translated_text = translate_text(run.text)
                               run.text = translated_text
                           set_arabic_format(paragraph)


       # Translate notes section in slides
       if slide.has_notes_slide:
           notes_slide = slide.notes_slide
           notes_text_frame = notes_slide.notes_text_frame
           if notes_text_frame:
               for paragraph in notes_text_frame.paragraphs:
                   for run in paragraph.runs:
                       translated_text = translate_text(run.text)
                       run.text = translated_text
                   set_arabic_format(paragraph)


   prs.save(output_file)


def translate_presentations_in_folder(input_folder, output_folder):
   if not os.path.exists(output_folder):
       os.makedirs(output_folder)


   for filename in os.listdir(input_folder):
       if filename.endswith('.pptx'):
           input_path = os.path.join(input_folder, filename)
           output_path = os.path.join(output_folder, f"translated_{filename}")
           translate_presentation(input_path, output_path)
           print(f"Translated {filename} and saved to {output_path}")


# Specify the input and output folders
input_folder = '/..'
output_folder = '/..'


translate_presentations_in_folder(input_folder, output_folder)


print("Translation of all presentations completed.")
